"""Deterministic OCR and generated parser samples for Phase 05 tests."""

from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO

from docx import Document
from openpyxl import Workbook
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches
from reportlab.pdfgen import canvas  # type: ignore[import-untyped]

from ragflow_agent.knowledge.domain.chunk import BoundingBox, CoordinateSpace
from ragflow_agent.knowledge.ports.ocr import OcrResult, OcrWord
from ragflow_agent.knowledge.ports.parsing import ParseRequest

DOCX_MEDIA_TYPE = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
PPTX_MEDIA_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
XLSX_MEDIA_TYPE = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"


@dataclass(frozen=True, slots=True)
class FormatSample:
    """One generated, source-controlled test recipe and expected output shape."""

    name: str
    media_type: str
    payload: bytes
    expected_kinds: frozenset[str]


class StaticOcrEngine:
    """Return stable geometry without pretending to test Tesseract itself."""

    def __init__(self, *, text: str = "alarm reset inspection") -> None:
        self._text = text
        self.calls: list[str] = []

    def recognize(self, image: object, *, language: str) -> OcrResult:
        del image
        self.calls.append(language)
        words = tuple(
            OcrWord(
                text=word,
                confidence=0.95,
                order=index,
                bounding_box=BoundingBox(
                    x0=float(10 + index * 80),
                    y0=12,
                    x1=float(70 + index * 80),
                    y1=36,
                    coordinate_space=CoordinateSpace.PIXELS,
                ),
            )
            for index, word in enumerate(self._text.split())
        )
        return OcrResult(
            engine_name="static-test-ocr",
            engine_version="1",
            language=language,
            words=words,
        )

    def available_languages(self) -> frozenset[str]:
        return frozenset({"eng", "chi_sim"})


def parse_request(sample: FormatSample) -> ParseRequest:
    """Build the stable source identity shared by direct parser tests."""
    return ParseRequest(
        tenant_id="tenant-a",
        knowledge_base_id="kb-a",
        document_id=f"doc-{sample.name}",
        document_version_id=f"version-{sample.name}",
        object_key=f"fixtures/{sample.name}",
        media_type=sample.media_type,
        trace_id=f"trace-{sample.name}",
    )


def generated_format_samples() -> tuple[FormatSample, ...]:
    """Generate legal, small, deterministic Phase 05 inputs at test time."""
    return (
        FormatSample(
            name="manual.txt",
            media_type="text/plain",
            payload=b"Alarm recovery\n\nReset the controller.",
            expected_kinds=frozenset({"text"}),
        ),
        FormatSample(
            name="manual.md",
            media_type="text/markdown",
            payload=(
                b"# Alarm Recovery\n\n"
                b"Reset the controller.\n\n"
                b"- Inspect relay\n\n"
                b"| Step | Action |\n| --- | --- |\n| 1 | Isolate power |\n"
            ),
            expected_kinds=frozenset({"heading", "text", "list", "table"}),
        ),
        FormatSample(
            name="manual.html",
            media_type="text/html",
            payload=(
                b"<h1>Alarm Recovery</h1><p>Reset controller.</p>"
                b"<script>secret()</script><table><tr><th>Step</th><th>Action</th></tr>"
                b"<tr><td>1</td><td>Inspect relay</td></tr></table>"
                b"<img src='relay.png' alt='Relay diagram'>"
            ),
            expected_kinds=frozenset({"heading", "text", "table", "image"}),
        ),
        FormatSample(
            name="manual.docx",
            media_type=DOCX_MEDIA_TYPE,
            payload=_docx_bytes(),
            expected_kinds=frozenset({"heading", "text", "table", "image"}),
        ),
        FormatSample(
            name="manual.pptx",
            media_type=PPTX_MEDIA_TYPE,
            payload=_pptx_bytes(),
            expected_kinds=frozenset({"heading", "text", "table", "image"}),
        ),
        FormatSample(
            name="alarms.xlsx",
            media_type=XLSX_MEDIA_TYPE,
            payload=_xlsx_bytes(),
            expected_kinds=frozenset({"table"}),
        ),
        FormatSample(
            name="manual.pdf",
            media_type="application/pdf",
            payload=_pdf_bytes(),
            expected_kinds=frozenset({"text", "table"}),
        ),
        FormatSample(
            name="alarm.png",
            media_type="image/png",
            payload=_image_bytes(),
            expected_kinds=frozenset({"image", "text"}),
        ),
    )


def blank_pdf_bytes() -> bytes:
    """Generate one blank page that must use the PDF OCR fallback."""
    stream = BytesIO()
    document = canvas.Canvas(stream, pagesize=(612, 792))
    document.showPage()
    document.save()
    return stream.getvalue()


def _docx_bytes() -> bytes:
    document = Document()
    document.add_heading("Alarm Recovery", level=1)
    document.add_paragraph("Reset the controller and inspect the relay.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Step"
    table.cell(0, 1).text = "Action"
    table.cell(1, 0).text = "1"
    table.cell(1, 1).text = "Isolate power"
    document.add_picture(BytesIO(_image_bytes()), width=Inches(1))
    stream = BytesIO()
    document.save(stream)
    return stream.getvalue()


def _pptx_bytes() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Alarm Recovery"
    placeholders = slide.placeholders
    placeholders[1].text = "Reset the controller."
    table = slide.shapes.add_table(
        2,
        2,
        Inches(1),
        Inches(4),
        Inches(5),
        Inches(1),
    ).table
    table.cell(0, 0).text = "Step"
    table.cell(0, 1).text = "Action"
    table.cell(1, 0).text = "1"
    table.cell(1, 1).text = "Inspect relay"
    slide.shapes.add_picture(
        BytesIO(_image_bytes()),
        Inches(6),
        Inches(4),
        width=Inches(1),
    )
    stream = BytesIO()
    presentation.save(stream)
    return stream.getvalue()


def _xlsx_bytes() -> bytes:
    workbook = Workbook()
    worksheet = workbook.active
    assert worksheet is not None
    worksheet.title = "Alarms"
    worksheet.append(["Code", "Action", "Priority"])
    worksheet.append(["A-100", "Reset controller", 1])
    worksheet.append(["A-200", '=CONCAT("Inspect", " relay")', 2])
    worksheet.merge_cells("A5:B5")
    worksheet["A5"] = "Approved procedure"
    second = workbook.create_sheet("Assets")
    second.append(["Asset", "Location"])
    second.append(["Relay-01", "Cabinet-A"])
    stream = BytesIO()
    workbook.save(stream)
    workbook.close()
    return stream.getvalue()


def _pdf_bytes() -> bytes:
    stream = BytesIO()
    document = canvas.Canvas(stream, pagesize=(612, 792))
    document.drawString(72, 740, "Alarm Recovery")
    document.drawString(72, 716, "Reset the controller and inspect the relay.")
    document.drawString(72, 690, "LEFT COLUMN first procedure")
    document.drawString(360, 690, "RIGHT COLUMN later evidence")
    left, top, cell_width, cell_height = 72, 660, 160, 24
    for row in range(3):
        y = top - row * cell_height
        document.line(left, y, left + 2 * cell_width, y)
    for column in range(3):
        x = left + column * cell_width
        document.line(x, top, x, top - 2 * cell_height)
    document.drawString(left + 4, top - 17, "Step")
    document.drawString(left + cell_width + 4, top - 17, "Action")
    document.drawString(left + 4, top - cell_height - 17, "1")
    document.drawString(
        left + cell_width + 4,
        top - cell_height - 17,
        "Isolate power",
    )
    document.save()
    return stream.getvalue()


def _image_bytes() -> bytes:
    image = Image.new("RGB", (480, 120), "white")
    ImageDraw.Draw(image).text((20, 40), "ALARM RESET INSPECTION", fill="black")
    stream = BytesIO()
    image.save(stream, format="PNG")
    return stream.getvalue()
