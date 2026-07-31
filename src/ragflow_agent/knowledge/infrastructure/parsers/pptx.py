"""Independent PPTX parser retaining slide geometry and source order."""

from __future__ import annotations

from io import BytesIO

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ragflow_agent.knowledge.domain.chunk import (
    BlockKind,
    BoundingBox,
    CoordinateSpace,
    ImageReference,
    ParsedBlock,
    ParseWarning,
    TableMetadata,
)
from ragflow_agent.knowledge.domain.errors import KnowledgeConflictError
from ragflow_agent.knowledge.infrastructure.parsers.security import validate_zip_package
from ragflow_agent.knowledge.ports.parsing import (
    ParsedPayload,
    ParserCapability,
    ParseRequest,
)


class PptxBinaryParser:
    """Parse slide titles, text, tables, and pictures without notes claims."""

    capability = ParserCapability(
        parser_id="independent-pptx",
        parser_version="2",
        media_types=frozenset(
            {
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            }
        ),
        extensions=frozenset({".pptx"}),
        default_chunk_strategy="manual",
    )

    def __init__(
        self,
        *,
        max_entries: int,
        max_uncompressed_bytes: int,
        max_compression_ratio: float,
    ) -> None:
        self._max_entries = max_entries
        self._max_uncompressed_bytes = max_uncompressed_bytes
        self._max_compression_ratio = max_compression_ratio

    def parse_bytes(
        self,
        payload: bytes,
        request: ParseRequest,
        *,
        ocr_language: str,
    ) -> ParsedPayload:
        del ocr_language
        validate_zip_package(
            payload,
            max_entries=self._max_entries,
            max_uncompressed_bytes=self._max_uncompressed_bytes,
            max_compression_ratio=self._max_compression_ratio,
        )
        try:
            presentation = Presentation(BytesIO(payload))
        except Exception as error:
            raise KnowledgeConflictError(
                "PPTX presentation is invalid",
                error_code="parser_pptx_invalid",
            ) from error
        blocks: list[ParsedBlock] = []
        warnings: list[ParseWarning] = []
        slide_width = max(
            int(presentation.slide_width) if presentation.slide_width is not None else 1,
            1,
        )
        slide_height = max(
            int(presentation.slide_height) if presentation.slide_height is not None else 1,
            1,
        )
        for slide_number, slide in enumerate(presentation.slides, start=1):
            heading_path: tuple[str, ...] = ()
            title_shape = slide.shapes.title
            for shape_index, shape in enumerate(slide.shapes):
                left = int(shape.left or 0)
                top = int(shape.top or 0)
                width = int(shape.width or 0)
                height = int(shape.height or 0)
                x0 = max(0.0, min(1.0, left / slide_width))
                y0 = max(0.0, min(1.0, top / slide_height))
                x1 = max(0.0, min(1.0, (left + width) / slide_width))
                y1 = max(0.0, min(1.0, (top + height) / slide_height))
                box = (
                    BoundingBox(
                        x0=x0,
                        y0=y0,
                        x1=x1,
                        y1=y1,
                        coordinate_space=CoordinateSpace.NORMALIZED,
                    )
                    if x1 > x0 and y1 > y0
                    else None
                )
                if getattr(shape, "has_table", False):
                    rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                    rows = [row for row in rows if any(row)]
                    if rows:
                        blocks.append(
                            ParsedBlock(
                                id=f"pptx_{len(blocks)}",
                                kind=BlockKind.TABLE,
                                order=len(blocks),
                                text="\n".join("\t".join(row) for row in rows),
                                page_number=slide_number,
                                bounding_box=box,
                                heading_path=heading_path,
                                table=TableMetadata(
                                    rows=len(rows),
                                    columns=max(len(row) for row in rows),
                                    has_header=True,
                                ),
                            )
                        )
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    blocks.append(
                        ParsedBlock(
                            id=f"pptx_{len(blocks)}",
                            kind=BlockKind.IMAGE,
                            order=len(blocks),
                            text="embedded PPTX image",
                            page_number=slide_number,
                            bounding_box=box,
                            heading_path=heading_path,
                            image=ImageReference(
                                object_key=request.object_key,
                                embedded_path=(f"slide-{slide_number}-picture-{shape_index}"),
                                media_type=str(shape.image.content_type),
                            ),
                        )
                    )
                elif getattr(shape, "has_text_frame", False):
                    value = str(shape.text).strip()
                    if not value:
                        continue
                    is_title = title_shape is not None and shape.shape_id == title_shape.shape_id
                    kind = BlockKind.HEADING if is_title else BlockKind.TEXT
                    if kind is BlockKind.HEADING:
                        heading_path = (value,)
                    blocks.append(
                        ParsedBlock(
                            id=f"pptx_{len(blocks)}",
                            kind=kind,
                            order=len(blocks),
                            text=value,
                            page_number=slide_number,
                            bounding_box=box,
                            heading_path=heading_path,
                        )
                    )
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    warnings.append(
                        ParseWarning(
                            code="pptx_notes_unsupported",
                            message="speaker notes are present but not indexed",
                            page_number=slide_number,
                        )
                    )
        return ParsedPayload(blocks=tuple(blocks), warnings=tuple(warnings))
